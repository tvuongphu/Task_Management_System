"""
Creates a PowerPoint presentation summarizing the Task Management System assignment.
Styled to match assignment-summary.html: dark background, cyan titles, light text.
Run: python scripts/create_presentation.py
Output: docs/Task-Management-Assignment-Summary.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Match assignment-summary.html theme
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)      # #1a1a2e
TITLE_CYAN = RGBColor(0x00, 0xD4, 0xFF)  # #00d4ff
TEXT_LIGHT = RGBColor(0xEE, 0xEE, 0xEE)  # #eee
SUBTITLE_MUTED = RGBColor(0xAA, 0xAA, 0xAA)  # #aaa


def _apply_dark_background(slide):
    """Apply dark background to slide (matches HTML #1a1a2e)."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _apply_dark_background(slide)
    slide.shapes.title.text = title
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = TITLE_CYAN
        run.font.size = Pt(40)
    if subtitle:
        sub_ph = slide.placeholders[1]
        sub_ph.text = subtitle
        for run in sub_ph.text_frame.paragraphs[0].runs:
            run.font.color.rgb = SUBTITLE_MUTED
            run.font.size = Pt(22)


def add_content_slide(prs, title, intro, bullet_points):
    """Add a content slide with optional intro paragraph and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _apply_dark_background(slide)

    # Title
    slide.shapes.title.text = title
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = TITLE_CYAN
        run.font.size = Pt(28)

    # Body
    body = slide.placeholders[1].text_frame
    body.clear()

    lines = []
    if intro:
        lines.append(intro)
    lines.extend(bullet_points)

    for i, text in enumerate(lines):
        p = body.add_paragraph()
        p.text = "▸ " + text.strip() if i > 0 else text.strip()  # Bullet for content lines
        p.level = 0
        p.space_after = Pt(6)
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = SUBTITLE_MUTED if i == 0 else TEXT_LIGHT


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Task Management System",
        "Assignment Summary — Architecture, Cloud, API, Testing & Process",
    )

    # Slide 2: Overview
    add_content_slide(
        prs,
        "Overview — Problem Statement",
        "Enterprise task platform: internal staff and external partners work together on tasks.",
        [
            "Create, assign, track tasks; comments and file attachments",
            "Who can do what: Azure AD sign-in + role-based access (RBAC)",
            "API for integrations (e.g. other systems) + web UI for people",
            "Scale: 5K–50K+ users, 100K–1M+ tasks per tenant",
            "Start simple (MVP), grow toward event sourcing and notifications",
        ],
    )

    # Slide 3: Deliverable 1 — Architecture
    add_content_slide(
        prs,
        "Deliverable 1: Solution Architecture",
        "We chose macroservices: one main app (core API) + one small service (notifications). Not full microservices (too much overhead), not one blob (notifications would slow the core).",
        [
            "Core monolith: tasks, comments, attachments, auth — one deployable",
            "Notification microservice: email, SMS, Teams — separate, scales independently",
            "Database: PostgreSQL — event store (every change saved) + projections (fast reads)",
            "Service Bus: core publishes events → notification service consumes — no tight coupling",
            "Identity: Azure AD (SSO, B2B for external partners)",
            "SignalR: real-time push — someone just edited this task",
        ],
    )

    # Slide 4: Architecture — Auth, Scale, Observability
    add_content_slide(
        prs,
        "Architecture — Auth, Scale, Observability",
        "Who gets in, who can do what, how we scale and monitor.",
        [
            "Auth: Azure AD SSO (single sign-on), JWT tokens; OAuth2 for app-to-app",
            "RBAC: roles (admin, member, viewer) — who can create, assign, delete tasks",
            "Scale: stateless API → add more instances; read replicas for queries; Redis for cache",
            "Observability: logs, traces, dashboards — Azure Monitor, Application Insights",
            "RED metrics: request rate, error rate, latency — SLO targets (e.g. p95 < 500 ms)",
        ],
    )

    # Slide 5: Deliverable 2 — Cloud
    add_content_slide(
        prs,
        "Deliverable 2: Cloud Setup (Azure)",
        "All services run on Azure — no on-prem servers.",
        [
            "Compute: App Service — hosts core API and notification service (separate plans for isolation)",
            "Data: PostgreSQL (DB), Redis (cache), Blob Storage (file attachments)",
            "Messaging: Service Bus (queues for notifications), SignalR Service (real-time push)",
            "Identity: Azure AD; email/SMS via Azure Communication Services",
            "Edge: Front Door or App Gateway — load balancing, SSL, WAF (DDoS protection)",
            "Secrets: Key Vault; deployments: Azure DevOps Pipelines",
        ],
    )

    # Slide 6: Cloud — Trade-offs
    add_content_slide(
        prs,
        "Cloud — Trade-offs & Cost",
        "Decisions that affect performance, cost, and resilience.",
        [
            "Core and notification on separate App Service plans — one can scale without affecting the other",
            "Cost: Start ~$120–250/mo (Dev/test) → Growth ~$1,550–2,600/mo (prod-scale)",
            "Front Door (global CDN) vs App Gateway (regional) — pick based on user geography",
            "Service Bus: Basic (simple queues) → Standard (topics) → Premium (VNet, compliance)",
            "Backup: Geo-redundant PostgreSQL and Blob — survive region outage",
        ],
    )

    # Slide 7: Deliverable 3 — API
    add_content_slide(
        prs,
        "Deliverable 3: API Design",
        "REST API for web UI and external systems. Base path: /v1/tasks",
        [
            "REST over JSON — POST (create), GET (list/search + get one), PATCH (update), DELETE",
            "Related: /tasks/{id}/comments, /tasks/{id}/attachments",
            "GET /tasks/{id}/history — full audit trail (who changed what, when)",
            "Concurrency: If-Match header or version in body — avoid overwriting someone else's edit",
            "Framework: ASP.NET Core; docs: OpenAPI/Swagger; testing: Postman",
        ],
    )

    # Slide 8: Deliverable 4 — Testing
    add_content_slide(
        prs,
        "Deliverable 4: Testing & QA",
        "Automate tests so QA focuses on exploratory work, not repetitive checks.",
        [
            "Pipeline: Build → Test + Security scan → Deploy to Dev → Staging → Prod (with approval)",
            "Unit: xUnit — test logic in isolation (NSubstitute for mocks)",
            "Integration: WebApplicationFactory + Testcontainers — real DB, test full request flow",
            "E2E: Playwright — browser tests; Contract: OpenAPI/Pact — API consumers stay aligned",
            "QA shift-left: involved in refinement and acceptance criteria before dev starts",
        ],
    )

    # Slide 9: Deliverable 5 — Process
    add_content_slide(
        prs,
        "Deliverable 5: Development Process",
        "How we work: process, team split, and how to speed up without overtime.",
        [
            "Team: 1 leader (tech lead/scrum master), 7 devs, 3 QA — 10 total. Allocation by phase in doc §1.4",
            "Scrum (2-week sprints), trunk-based branching. DoR = story ready before pull; DoD = tests + QA sign-off",
            "Speed 2–3x without overtime: parallel delivery, CI/CD, automation, feature flags",
            "Parallel working: Week 1 — agree event contract. Team A (4 devs) Phase 1 core; Team B (3 devs) Notification + SignalR. QA split per stream. Integrate when core write path ready.",
            "Metrics: DORA (deploy freq, lead time, MTTR). Risks: scope creep, key-person dependency — hard capacity, cross-training",
        ],
    )

    # Slide 10: Roadmap
    add_content_slide(
        prs,
        "Project Roadmap (Weeks 1–12)",
        "Phased delivery — each phase is shippable.",
        [
            "Phase 1 (weeks 1–4): Event store, task CRUD, comments, attachments, auth, CI/CD. Core API works; optional minimal UI.",
            "Phase 2 (weeks 5–8): SignalR (real-time), Notification microservice, Service Bus. Can run in parallel with Phase 1 if event contract agreed early.",
            "Phase 3 (weeks 9–12): Full SPA, scheduler jobs, production deployment pipeline, monitoring.",
            "Phase 4: UAT, go-live, iterate on feedback",
        ],
    )

    # Slide 11: Summary
    add_content_slide(
        prs,
        "Summary",
        "One-line takeaway per deliverable.",
        [
            "Architecture: Macroservices — core monolith + Notification microservice; event sourcing; Azure Service Bus + SignalR",
            "Cloud: Azure-native — App Service, PostgreSQL, Redis, Blob, Service Bus, SignalR, Key Vault",
            "API: REST, /tasks CRUD + comments + attachments + history; ASP.NET Core",
            "Testing: CI/CD pipeline; unit, integration, contract, E2E (Playwright); QA shift-left",
            "Process: Team 7 devs + 3 QA; parallel delivery (event contract → Team A core, Team B Notification+SignalR); Scrum, DORA metrics",
        ],
    )

    # Save
    output_path = "docs/Task-Management-Assignment-Summary.pptx"
    prs.save(output_path)
    print(f"Created: {output_path}")


if __name__ == "__main__":
    main()
